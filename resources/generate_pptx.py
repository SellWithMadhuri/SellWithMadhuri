import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_dholera_presentation():
    prs = Presentation()
    # Set 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Color Palette
    DARK_TEAL = RGBColor(11, 43, 38)     # #0B2B26
    MID_TEAL = RGBColor(15, 76, 58)      # #0F4C3A
    ACCENT_GOLD = RGBColor(234, 179, 8)   # #EAB308
    ACCENT_GREEN = RGBColor(34, 197, 94)  # #22C55E
    LIGHT_BG = RGBColor(240, 253, 244)   # #F0FDF4
    TEXT_DARK = RGBColor(15, 23, 42)     # #0F172A
    TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
    WHITE = RGBColor(255, 255, 255)

    slides_data = [
        # Slide 1
        {"type": "title", "title": "DHOLERA SMART CITY", "subtitle": "India's 1st Planned Greenfield Smart Megacity\nSanctuary Home Plotting Project | ABDB Realtech"},
        
        # Slide 2
        {"type": "content", "title": "1. Executive Summary - The Dholera Horizon", 
         "bullets": ["India's premier Greenfield Smart City under the Delhi-Mumbai Industrial Corridor (DMIC).",
                     "Total Planned Region: 920 Sq. Km (22 Town Planning schemes).",
                     "Sanctuary Home Project: Premier title-clear residential plotting starting at ₹6,250 / sq. yd.",
                     "Designed for 2.0 Million citizens with 800,000+ upcoming high-value jobs."]},

        # Slide 3
        {"type": "content", "title": "2. Grand Scale - Larger Than Singapore", 
         "bullets": ["Total Special Investment Region: 920 Sq. Km (91,970 Hectares).",
                     "Singapore Comparison: Singapore spans 728 Sq. Km vs Dholera's 920 Sq. Km.",
                     "Activation Area: 22.5 Sq. Km fully developed with trunk infrastructure.",
                     "Strategic hub positioning India as a global manufacturing superpower."]},

        # Slide 4
        {"type": "content", "title": "3. Strategic Location & DMIC Corridor", 
         "bullets": ["Situated in Ahmedabad district, Gujarat, along the Gulf of Khambhat.",
                     "Core node of the $100 Billion Delhi-Mumbai Industrial Corridor (DMIC).",
                     "Direct connectivity to major industrial ports, highways, and rail nodes.",
                     "High Access Corridor linking Dholera seamlessly with Ahmedabad and Bhavnagar."]},

        # Slide 5
        {"type": "content", "title": "4. Master Plan & Smart Zoning", 
         "bullets": ["Eco-Conscious Ratio: 60% Greenery & Open Spaces vs 40% Development.",
                     "Zoning Breakdown: Industrial, Knowledge & IT, Public Facilities, Residential, Logistics.",
                     "High Access Corridor & City Centre Zone for corporate headquarters.",
                     "Sustainable urban design preventing city congestion and pollution."]},

        # Slide 6
        {"type": "content", "title": "5. Demographic Target - 2.0 Million Residents", 
         "bullets": ["Target Population: Planned for 2.0 Million residents by 2040.",
                     "High Quality of Life: Walk-to-work culture, smart housing, and pristine environment.",
                     "World-class civic infrastructure: 24/7 power, SCADA water supply, advanced waste recycling.",
                     "Designed to attract top international tech talent and executive families."]},

        # Slide 7
        {"type": "content", "title": "6. Employment Engine - 800,000+ Jobs", 
         "bullets": ["Projected Job Creation: 8 Lakh to 10 Lakh industrial & tech jobs.",
                     "Anchor Heavy Industries: Defense, Aerospace, Electronics, Automotive, Solar Energy.",
                     "Knowledge & IT Hubs: R&D parks, software campuses, and skill universities.",
                     "Tremendous housing demand creation for incoming corporate workforce."]},

        # Slide 8
        {"type": "content", "title": "7. India's Semiconductor Hub (TATA Fab)", 
         "bullets": ["TATA Semiconductor Plant: ₹91,000 Crore investment in Dholera SIR.",
                     "Commercial Production: First chip production targeted by 2026.",
                     "Ecosystem Impact: 50+ ancillary component manufacturers setting up facilities.",
                     "Massive boost to local land values and immediate residential rental demand."]},

        # Slide 9
        {"type": "content", "title": "8. The 5 Mega Connectivity Pillars", 
         "bullets": ["1. Dholera International Airport (Under construction, Dec 2026 flight).",
                     "2. Ahmedabad-Dholera Expressway (4-Lane High Speed Corridor).",
                     "3. High-Speed Metro Rail Project.",
                     "4. Deepwater Seaport Project.",
                     "5. Western Dedicated Freight Corridor (DFC)."]},

        # Slide 10
        {"type": "content", "title": "9. Dholera International Airport (Dec 2026)", 
         "bullets": ["Budget Allocation: ₹5,083 Crore approved by Ministry of Civil Aviation.",
                     "Phase 1 Cargo Flight Target: Expected operational by December 2026.",
                     "Location: Navagam village, serving both passenger passenger and heavy cargo logistics.",
                     "Distance from Sanctuary Home: 30 minutes drive."]},

        # Slide 11
        {"type": "content", "title": "10. Ahmedabad-Dholera Expressway", 
         "bullets": ["Budget Allocation: ₹3,200 to ₹4,000 Crore by NHAI.",
                     "High-Speed Expressway: Cuts travel time between Ahmedabad & Dholera to 45 minutes.",
                     "Current Progress: Operational timeline targeted for 2025-2026.",
                     "Proximity to Project: Sanctuary Home is just 4 to 5 minutes from Expressway access."]},

        # Slide 12
        {"type": "content", "title": "11. Dholera Metro Rail Project", 
         "bullets": ["Approved Budget: ₹12,500 – ₹13,000 Crore by Central Government.",
                     "Connecting Hubs: Links Ahmedabad Metro directly with Dholera SIR Activation Area.",
                     "Target Operations: Expected to commence around 2027.",
                     "Elevated rapid transit providing effortless commute for residents."]},

        # Slide 13
        {"type": "content", "title": "12. Dholera Deepwater Seaport", 
         "bullets": ["Approved Budget: ₹110 Crore sanctioned by Gujarat Maritime Board.",
                     "Strategic Function: Facilitates international maritime cargo and export-import trade.",
                     "Operational Timeline: Expected completion between 2027 – 2028.",
                     "Gives industrial tenants direct access to global shipping lanes."]},

        # Slide 14
        {"type": "content", "title": "13. Dedicated Freight Corridor (DFC)", 
         "bullets": ["Approved Budget: ₹466 Crore by Ministry of Railways.",
                     "High Capacity Logistics: Connects Dholera to Northern & Western Industrial Belts.",
                     "Target Completion: Set to be fully operational by late 2025.",
                     "Drastically reduces freight transport costs for manufacturing units."]},

        # Slide 15
        {"type": "content", "title": "14. ABCD Building - Brain & Heart of SIR", 
         "bullets": ["ABCD: Activation Area Control, Command & Development Building.",
                     "Functions as the Central City Integrated Operations Centre (CIOC).",
                     "Single-Window Clearance System: Streamlines project approvals for investors.",
                     "Completed & Operational: 1st fully functional government administrative hub in SIR."]},

        # Slide 16
        {"type": "content", "title": "15. Smart Utilities & IoT Network", 
         "bullets": ["100% Underground Utilities: Zero overhead cables, smart grid power distribution.",
                     "SCADA Water Management: Desalination and 100% recycled wastewater lines.",
                     "City-Wide Sensor Network: AI traffic control, automatic fault detection, smart streetlights.",
                     "Safety Sensors: 24/7 gas and fire leak detection across homes and workplaces."]},

        # Slide 17
        {"type": "content", "title": "16. Sanctuary Home - Premium Plotting Project", 
         "bullets": ["Developer: ABDB Realtech Pvt. Ltd. (Proven real estate track record).",
                     "Project Type: High-end gated residential plotting township.",
                     "Vision: A perfect beginning for a better tomorrow - luxury living + high ROI.",
                     "Starting Price: Highly attractive initial rate of ₹6,250 / sq. yd."]},

        # Slide 18
        {"type": "content", "title": "17. Strategic Location Advantage", 
         "bullets": ["100 to 200 Meters from SIR Boundary (High appreciation zone).",
                     "4 to 5 Minutes drive from Ahmedabad-Dholera Expressway.",
                     "2 Minutes from local School, Bank, Post Office & Kamiyala Temple.",
                     "20 Minutes to ABCD Command Centre & 30 Minutes to International Airport."]},

        # Slide 19
        {"type": "content", "title": "18. Approved Layout & Master Plan", 
         "bullets": ["Total Inventory: 56 Sanctioned Residential Plots.",
                     "Internal Infrastructure: Wide paved roads with street light poles.",
                     "Common Open Spaces: Dedicated COP-1, COP-2 & COP-3 green reserves.",
                     "Optimal Layout: Designed for maximum natural ventilation and Vastu harmony."]},

        # Slide 20
        {"type": "content", "title": "19. Plot Options & Size Flexibility", 
         "bullets": ["Standard Compact Plots: 143 sq. yd. to 191 sq. yd. (Ideal for modular villas).",
                     "Medium Family Plots: 250 sq. yd. to 350 sq. yd.",
                     "Premium Corner & Bungalow Plots: 426 sq. yd., 513 sq. yd. up to 758 sq. yd.",
                     "Flexibility: Freedom to design custom architectural villas."]},

        # Slide 21
        {"type": "content", "title": "20. Investment Pricing & Land Rates", 
         "bullets": ["Introductory Price: ₹6,250 per sq. yd.",
                     "Plot Investment Starting: ~₹15.6 Lakhs for 250 sq. yd. plot.",
                     "High Capital Growth: Land prices projected to 2.5x to 4x as airport opens.",
                     "Flexible payment plans & transparent pricing structure."]},

        # Slide 22
        {"type": "content", "title": "21. Premium Amenities for Quality Lifestyle", 
         "bullets": ["Gated Gentry Society with Guard Room & Barricades.",
                     "24x7 Security & CCTV Surveillance Network.",
                     "Lord Shiv Temple built within the society premises.",
                     "3 Landscaped Green Parks & Dedicated Kids Play Area."]},

        # Slide 23
        {"type": "content", "title": "22. Legal Security & 100% Clear Title", 
         "bullets": ["NA (Non-Agricultural) Approved land status.",
                     "NOC & Title Clear documentation from local authorities.",
                     "Immediate Registry with Official Land Mutation.",
                     "Zero legal risk: Complete peace of mind for buyers and investors."]},

        # Slide 24
        {"type": "content", "title": "23. Diamond Circle Development Node", 
         "bullets": ["Strategic Intersection: Key development node connecting internal mobility roads.",
                     "Budget Allocation: ₹100 – ₹150 Crore infrastructure project.",
                     "Operational Timeline: Targeted completion mid-2025.",
                     "Significantly enhances local traffic flow and boosts nearby property values."]},

        # Slide 25
        {"type": "content", "title": "24. Eco-Tourism & Surroundings", 
         "bullets": ["Blackbuck National Park (Velavadar): Famous tourist attraction nearby.",
                     "Upcoming Mega Sea Link Project: Just 15 minutes distance.",
                     "Pristine coastal breeze & pollution-free environment.",
                     "Combines urban smart living with serene eco-lifestyle."]},

        # Slide 26
        {"type": "content", "title": "25. Land Appreciation Projections (2026-2030)", 
         "bullets": ["Phase 1 Trigger (Expressway 2025): Estimated 30%-40% value surge.",
                     "Phase 2 Trigger (Airport Cargo 2026): Estimated 60%-80% surge.",
                     "Phase 3 Trigger (TATA Fab Production): Industrial workforce demand boom.",
                     "5-Year Projection: High probability of 3x to 5x returns on initial land investment."]},

        # Slide 27
        {"type": "content", "title": "26. Investor Case Study - Early Mover Advantage", 
         "bullets": ["Historical Trend: Early investors in GIFT City & Gurgaon saw 10x gains.",
                     "Dholera Advantage: Government-backed infrastructure built *before* population arrival.",
                     "Low Entry Barrier: Land starting at ₹6,250/sq.yd compared to developed cities.",
                     "Ideal asset class for long-term wealth creation and family legacy."]},

        # Slide 28
        {"type": "content", "title": "27. Easy 4-Step Purchase Process", 
         "bullets": ["Step 1: Select plot location & size from layout map.",
                     "Step 2: Complete booking registration with token amount.",
                     "Step 3: Verification of NA/NOC title clear documents.",
                     "Step 4: Registry execution with official mutation certificate."]},

        # Slide 29
        {"type": "content", "title": "28. Free Guided Site Visit Program", 
         "bullets": ["Complimentary Pick & Drop service for interested buyers.",
                     "Guided tour of Dholera SIR Activation Area, ABCD Building & Expressway.",
                     "On-site physical plot inspection at Sanctuary Home.",
                     "Direct consultation with ABDB Realtech project advisors."]},

        # Slide 30
        {"type": "content", "title": "29. Why Choose ABDB Realtech", 
         "bullets": ["Proven Real Estate Integrity: Transparent legal procedures.",
                     "Customer First Policy: Dedicated post-purchase mutation & registry assistance.",
                     "Corporate Presence: Office in Noida & local team in Dholera.",
                     "Track record of identifying high-growth corridors."]},

        # Slide 31 (Closing)
        {"type": "title", "title": "BUILD YOUR DREAM - LIVE YOUR FUTURE", 
         "subtitle": "Book Your Sanctuary Home Plot Today!\n\nDeveloper: ABDB Realtech Pvt. Ltd.\nWebsite: www.abdbrealtech.com | Phone: +91 92113 13566\nEmail: Md.abdbrealtech@gmail.com"}
    ]

    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Background shape fill
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.line.fill.background()
        
        if slide_info["type"] == "title":
            bg.fill.solid()
            bg.fill.fore_color.rgb = DARK_TEAL
            
            # Title text frame
            tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
            tf = tx_box.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = slide_info["title"]
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = ACCENT_GOLD
            p.alignment = PP_ALIGN.CENTER
            
            p2 = tf.add_paragraph()
            p2.text = "\n" + slide_info["subtitle"]
            p2.font.size = Pt(24)
            p2.font.color.rgb = WHITE
            p2.alignment = PP_ALIGN.CENTER

        else:
            bg.fill.solid()
            bg.fill.fore_color.rgb = LIGHT_BG
            
            # Header banner
            header_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
            header_rect.fill.solid()
            header_rect.fill.fore_color.rgb = DARK_TEAL
            header_rect.line.fill.background()
            
            # Header text
            tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8))
            tf = tx_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_info["title"]
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = WHITE
            
            # Content Card Container
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = MID_TEAL
            
            # Bullet text frame
            c_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(11.0), Inches(4.6))
            ctf = c_box.text_frame
            ctf.word_wrap = True
            
            for idx, bullet in enumerate(slide_info["bullets"]):
                p = ctf.paragraphs[0] if idx == 0 else ctf.add_paragraph()
                p.text = "• " + bullet
                p.font.size = Pt(22)
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(18)

    output_path = r"d:\Dholera\Dholera_Smart_City_30_Slide_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_dholera_presentation()
